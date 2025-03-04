from pptx import Presentation
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from datetime import datetime
from typing import Union, Tuple, List, Dict, BinaryIO
from sqlalchemy.orm import Session
from app.models.models import SlideMetadata
import shutil
from pathlib import Path
from pptx.shapes.placeholder import SlidePlaceholder
import json
from app.utils.openai import get_embedding



async def process_powerpoint_repository(
    pptx_source: Union[str, BinaryIO, bytes],
    db: Session,
    source_type: str = "file_path"
) -> Tuple[str, List[SlideMetadata]]:
    """
    Process a PowerPoint file and create SlideMetadata objects.
    Stores the PowerPoint file in slides_storage directory.
    
    Args:
        pptx_source: Can be either:
            - file path (str)
            - file-like object (BinaryIO from upload)
            - bytes (from MS Graph API)
        db: Session: SQLAlchemy database session
        source_type: One of "file_path", "upload", or "ms_graph"
    
    Returns:
        Tuple containing:
            - storage_path: Path where the presentation was stored
            - List of SlideMetadata objects
    """
    # Create storage directory if it doesn't exist
    storage_dir = Path("slides_repository")
    storage_dir.mkdir(exist_ok=True)
    
    # Generate unique filename using timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    storage_path = storage_dir / f"presentation_{timestamp}.pptx"
    
    # Handle different input types and save to storage
    if source_type == "file_path":
        shutil.copy2(pptx_source, storage_path)
        presentation = Presentation(pptx_source)
    elif source_type == "upload":
        with open(storage_path, "wb") as f:
            if hasattr(pptx_source, "seek"):
                pptx_source.seek(0)
            f.write(pptx_source.read())
        presentation = Presentation(storage_path)
    elif source_type == "ms_graph":
        with open(storage_path, "wb") as f:
            f.write(pptx_source)
        presentation = Presentation(storage_path)
    else:
        raise ValueError("Invalid source_type. Must be 'file_path', 'upload', or 'ms_graph'")

    # Create metadata objects for each slide
    slide_metadata_objects = []
    
    for slide_idx, slide in enumerate(presentation.slides):
        semantic_content = {
            "title": _extract_slide_title(slide),
            "purpose": _infer_slide_purpose(slide),
            "category": _infer_slide_category(slide),
            "tags": _generate_slide_tags(slide)
        }
        stringified_metadata = json.dumps(semantic_content)
        embedding = await get_embedding(stringified_metadata)
        
        metadata = SlideMetadata(
            title=semantic_content["title"],
            category=semantic_content["category"],
            slide_type=_infer_slide_type(slide),
            purpose=semantic_content["purpose"],
            tags=semantic_content["tags"],
            audience=None,  # To be filled by user/AI later
            sales_stage=None,  # To be filled by user/AI later
            content_schema=_create_content_schema(slide),
            embedding=embedding
        )
        
        slide_metadata_objects.append(metadata)
    
    return str(storage_path), slide_metadata_objects

def _extract_slide_title(slide) -> str:
    """Extract the title from a slide"""
    for shape in slide.shapes:
        try:
            # First check if it's a placeholder and if it's a title
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                if shape.has_text_frame:
                    return shape.text_frame.text
        except (AttributeError, ValueError):
            # If it's not a placeholder, check if it might be a title shape by name
            if hasattr(shape, 'name') and 'Title' in shape.name:
                if shape.has_text_frame:
                    return shape.text_frame.text
    
    return "Untitled Slide"

def _infer_slide_category(slide) -> str:
    """Infer the category of the slide based on its content"""
    # Check for charts
    if any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in slide.shapes):
        return "data_visualization"
    
    # Check for tables
    if any(hasattr(shape, "table") for shape in slide.shapes):
        return "data_presentation"
    
    # Check for multiple images
    if sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) > 1:
        return "visual"
    
    # Default to content slide
    return "content"

def _infer_slide_type(slide) -> str:
    """Infer the specific type of slide based on its content and layout"""
    layout_name = slide.slide_layout.name.lower()
    
    # Check layout name first
    if "title" in layout_name:
        if "content" in layout_name:
            return "title_and_content"
        return "title_slide"
    
    # Check content types
    if any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in slide.shapes):
        return "chart_slide"
    
    if any(hasattr(shape, "table") for shape in slide.shapes):
        return "table_slide"
    
    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        return "image_slide"
    
    return "text_slide"

def _infer_slide_purpose(slide) -> str:
    """Infer the purpose of the slide based on its type and content"""
    slide_type = _infer_slide_type(slide)
    
    # Basic mapping of slide types to purposes
    purpose_map = {
        "title_slide": "To introduce the presentation",
        "chart_slide": "To visualize data or trends",
        "table_slide": "To present structured data",
        "image_slide": "To provide visual information",
        "text_slide": "To convey textual information"
    }
    
    return purpose_map.get(slide_type, "To present information")

def _generate_slide_tags(slide) -> List[str]:
    """Generate relevant tags for the slide"""
    tags = []
    
    # Add layout-based tag
    tags.append(slide.slide_layout.name.lower())
    
    # Add content-type based tags
    content_types = set()
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            content_types.add("chart")
        elif hasattr(shape, "table"):
            content_types.add("table")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            content_types.add("picture")
        elif shape.has_text_frame:
            content_types.add("text")
    
    tags.extend(list(content_types))
    
    return tags

#TODO: This is a placeholder function to create a content object/schema for the slides.
def _create_content_schema(slide) -> Dict:
    """
    Create comprehensive content object for a given slide.
    
    Args:
        slide (Slide): The slide to create a content object for
    
    Returns:
        dict: Dictionary containing the content object for the slide
    """
    # Get placeholder type names for better readability
    placeholder_types = {
        getattr(PP_PLACEHOLDER, attr): attr 
        for attr in dir(PP_PLACEHOLDER) 
        if not attr.startswith('__')
    }
    
    # Get shape type names
    shape_types = {
        getattr(MSO_SHAPE_TYPE, attr): attr 
        for attr in dir(MSO_SHAPE_TYPE) 
        if not attr.startswith('__')
    }
    
    # Basic slide information
    schema = {
        "slide_id": slide.slide_id,
        "layout_name": slide.slide_layout.name,
        "elements": []
    }
    
    # Process all shapes on the slide
    for shape in slide.shapes:
        # Base shape information
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "shape_type": shape_types.get(shape.shape_type, str(shape.shape_type)),
            "position": {
                "x": shape.left,
                "y": shape.top,
                "width": shape.width,
                "height": shape.height,
                "rotation": shape.rotation
            }
        }
        
        # Add placeholder information if applicable
        try:
            if shape.is_placeholder:
                element["is_placeholder"] = True
                placeholder_type = shape.placeholder_format.type
                element["placeholder_type"] = placeholder_types.get(placeholder_type, str(placeholder_type))
                element["placeholder_idx"] = shape.placeholder_format.idx
            else:
                element["is_placeholder"] = False
        except (AttributeError, ValueError):
            element["is_placeholder"] = False
        
        # Text content
        if shape.has_text_frame:
            element["content_type"] = "text"
            
            # Get text and formatting at paragraph level
            paragraphs_data = []
            for p in shape.text_frame.paragraphs:
                paragraph_data = {
                    "text": p.text,
                    "level": p.level,
                    "alignment": str(p.alignment) if hasattr(p, 'alignment') else None,
                    "runs": []
                }
                
                # Get formatting for each text run
                for run in p.runs:
                    run_data = {
                        "text": run.text,
                        "font": {
                            "name": run.font.name,
                            "size": run.font.size.pt if hasattr(run.font.size, 'pt') else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                            "color": str(run.font.color.rgb) if hasattr(run.font.color, 'rgb') and run.font.color.rgb else None
                        }
                    }
                    paragraph_data["runs"].append(run_data)
                
                paragraphs_data.append(paragraph_data)
            
            element["paragraphs"] = paragraphs_data
            element["has_text_linking"] = shape.text_frame.auto_size
            element["word_wrap"] = shape.text_frame.word_wrap
            element["vertical_anchor"] = str(shape.text_frame.vertical_anchor) if hasattr(shape.text_frame, 'vertical_anchor') else None
        
        # Table content
        elif hasattr(shape, 'table'):
            element["content_type"] = "table"
            table_data = []
            
            for r_idx, row in enumerate(shape.table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    cell_data = {
                        "text": cell.text,
                        "row_idx": r_idx,
                        "col_idx": c_idx,
                        "row_span": cell.span.row_span,
                        "col_span": cell.span.col_span,
                        "width": cell.width,
                        "height": cell.height
                    }
                    row_data.append(cell_data)
                table_data.append(row_data)
            
            element["table_data"] = table_data
            element["row_count"] = len(shape.table.rows)
            element["column_count"] = len(shape.table.columns)
        
        # Chart content
        elif hasattr(shape, 'chart'):
            element["content_type"] = "chart"
            element["chart_type"] = str(shape.chart.chart_type)
            
            # Extract categories and series
            if hasattr(shape.chart, 'plots') and shape.chart.plots:
                plot = shape.chart.plots[0]
                
                # Try to get categories
                categories = []
                if hasattr(plot, 'categories') and plot.categories:
                    for category in plot.categories:
                        categories.append(str(category.label))
                element["categories"] = categories
                
                # Try to get series
                series_data = []
                if hasattr(plot, 'series'):
                    for series in plot.series:
                        series_info = {
                            "name": series.name if hasattr(series, 'name') else "Unknown",
                            "values": list(series.values) if hasattr(series, 'values') else []
                        }
                        series_data.append(series_info)
                element["series"] = series_data
            
            element["has_title"] = shape.chart.has_title
            if shape.chart.has_title:
                element["title"] = shape.chart.chart_title.text_frame.text
        
        # Picture content
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["content_type"] = "picture"
            if hasattr(shape, 'image'):
                element["image_filename"] = shape.image.filename if hasattr(shape.image, 'filename') else None
                element["image_format"] = shape.image.ext if hasattr(shape.image, 'ext') else None
                element["content_type"] = f"image/{shape.image.ext}" if hasattr(shape.image, 'ext') else "image"
        
        # Shape content
        else:
            element["content_type"] = "shape"
            if hasattr(shape, 'fill'):
                element["fill_type"] = str(shape.fill.type) if hasattr(shape.fill, 'type') else None
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                    element["fill_color"] = str(shape.fill.fore_color.rgb) if hasattr(shape.fill.fore_color, 'rgb') else None
        
        # Add this element to the slide schema
        schema["elements"].append(element)
    
    # Add background information
    try:
        if hasattr(slide, 'background') and slide.background:
            if hasattr(slide.background, 'fill'):
                fill_type = type(slide.background.fill._fill).__name__
                schema["background_fill_type"] = fill_type
                
                # Only try to get fore_color if it's a solid fill
                if fill_type == "SolidFill" and hasattr(slide.background.fill, 'fore_color'):
                    if slide.background.fill.fore_color and hasattr(slide.background.fill.fore_color, 'rgb'):
                        schema["background_color"] = str(slide.background.fill.fore_color.rgb)
    except (AttributeError, TypeError):
        # If anything goes wrong with background processing, just skip it
        pass
    
    return schema



def construct_presentation(template_path, output_path, slide_data):
    """
    Copy selected slides from a template presentation, insert/modify content of the slides based the slide_data, and create a new presentation.
    
    Args:
        template_path (str): Path to the template presentation
        output_path (str): Path where the new presentation will be saved
        slide_data (list): List of dictionaries with slide_id and content to be updated
                          Example: [
                              {
                                  "slide_id": 5,  # slide number or some identifier
                                  "content": {
                                      "title": "New Title",
                                      "subtitle": "New Subtitle",
                                      "body": "New content for body placeholder"
                                  }
                              },
                              ...
                          ]
    """
    # Load the template presentation
    template_pres = Presentation(template_path)
    
    # Create a new presentation to hold the copied slides
    new_pres = Presentation()
    
    # Create a mapping of placeholder types to their names for easier identification
    placeholder_types = {getattr(PP_PLACEHOLDER, attr): attr 
                        for attr in dir(PP_PLACEHOLDER) 
                        if not attr.startswith('__')}
    
    # Map slide IDs to their indices for easier lookup
    # Note: If slide_id in your data is actual slide index, this step isn't necessary
    slide_map = {i: slide for i, slide in enumerate(template_pres.slides)}
    
    # Process each slide in the slide_data
    for slide_item in slide_data:
        slide_id = slide_item["slide_id"]
        content = slide_item["content"]
        
        # Get the template slide (assuming slide_id is the index; adjust if needed)
        if slide_id not in slide_map:
            print(f"Warning: Slide ID {slide_id} not found in template presentation")
            continue
            
        template_slide = slide_map[slide_id]
        
        # Create a blank slide in the new presentation with the same layout
        layout = template_slide.slide_layout
        new_slide = new_pres.slides.add_slide(layout)
        
        # Copy and modify each shape from the template slide
        for shape in template_slide.shapes:
            # If it's a placeholder that needs content update
            if hasattr(shape, "placeholder_format") and isinstance(shape, SlidePlaceholder):
                # Get placeholder type
                ph_type = shape.placeholder_format.type
                ph_type_name = placeholder_types.get(ph_type, "UNKNOWN")
                
                # Find the corresponding placeholder in the new slide
                for new_shape in new_slide.shapes:
                    if (hasattr(new_shape, "placeholder_format") and 
                        new_shape.placeholder_format.type == ph_type):
                        
                        # Update placeholder content if it exists in our content dictionary
                        # Convert placeholder type to lowercase for key matching
                        content_key = ph_type_name.lower()
                        
                        # Allow flexible key matching
                        matched_key = None
                        if content_key in content:
                            matched_key = content_key
                        elif shape.name.lower() in content:
                            matched_key = shape.name.lower()
                        elif "title" in content_key and "title" in content:
                            matched_key = "title"
                        elif "subtitle" in content_key and "subtitle" in content:
                            matched_key = "subtitle"
                        elif "body" in content_key and "body" in content:
                            matched_key = "body"
                        
                        if matched_key:
                            if hasattr(new_shape, "text_frame"):
                                new_shape.text_frame.text = content[matched_key]
                            elif hasattr(new_shape, "text"):
                                new_shape.text = content[matched_key]
            
            # Process other shape types as needed
            # (For images, tables, charts, etc., you would need more complex logic)
    
    # Save the new presentation
    new_pres.save(output_path)
    return new_pres



def duplicate_slide_object(slide):
    """
    Create a duplicate of a slide object
    
    Args:
        slide: The source Slide object to duplicate
        
    Returns:
        A new Slide object with duplicated content
    """
    # Get the presentation that contains the slide
    prs = slide.part.package.presentation
    
    # Create a temporary new slide with the same layout
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy slide content
    for shape in slide.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    # Copy slide properties
    if hasattr(slide, 'background') and hasattr(new_slide, 'background'):
        new_slide.background = slide.background
    
    # Return the new slide object
    return new_slide

    

def apply_modified_schema(slide, original_schema, modified_schema):
    """
    Apply only the modified parts of a schema to a slide
    
    Args:
        slide: The slide object to apply changes to (a duplicated slide)
        original_schema: Dictionary with the original slide schema
        modified_schema: Dictionary with the modified slide schema
        
    Returns:
        slide: The updated slide object
    """
    # We'll iterate through the elements in the modified schema
    # and update only what has changed compared to the original schema
    
    for mod_element in modified_schema["elements"]:
        # Find the corresponding element in the original schema
        orig_element = next(
            (e for e in original_schema["elements"] if e["id"] == mod_element["id"]), 
            None
        )
        
        # If this is a new element that wasn't in the original, we'd handle it differently
        # (that would require creating new shapes, which is more complex)
        if not orig_element:
            print(f"Warning: Element ID {mod_element['id']} not found in original schema")
            continue
        
        # Find the corresponding shape in the slide
        shape = None
        for s in slide.shapes:
            if s.shape_id == mod_element["id"]:
                shape = s
                break
        
        if not shape:
            print(f"Warning: Shape with ID {mod_element['id']} not found in slide")
            continue
        
        # Now we update only what has changed
        if "position" in mod_element:
            # Check if position attributes have changed
            mod_pos = mod_element["position"]
            orig_pos = orig_element["position"]
            
            if mod_pos.get("x") != orig_pos.get("x"):
                shape.left = mod_pos["x"]
            if mod_pos.get("y") != orig_pos.get("y"):
                shape.top = mod_pos["y"]
            if mod_pos.get("width") != orig_pos.get("width"):
                shape.width = mod_pos["width"]
            if mod_pos.get("height") != orig_pos.get("height"):
                shape.height = mod_pos["height"]
            if mod_pos.get("rotation") != orig_pos.get("rotation"):
                shape.rotation = mod_pos["rotation"]
        
        # Handle text content changes
        if mod_element.get("content_type") == "text" and shape.has_text_frame:
            mod_paras = mod_element.get("paragraphs", [])
            orig_paras = orig_element.get("paragraphs", [])
            
            # Only update if there are differences in paragraphs
            if mod_paras != orig_paras:
                # Clear existing text frame content
                for i in range(len(shape.text_frame.paragraphs)-1, 0, -1):
                    p = shape.text_frame.paragraphs[i]
                    p._p.getparent().remove(p._p)
                
                first_para = shape.text_frame.paragraphs[0]
                first_para.text = ""
                
                # Apply new text content
                for i, para_data in enumerate(mod_paras):
                    if i == 0:
                        para = first_para
                    else:
                        para = shape.text_frame.add_paragraph()
                    
                    # Apply paragraph-level formatting
                    if "level" in para_data:
                        para.level = para_data["level"]
                    
                    if "alignment" in para_data and para_data["alignment"]:
                        # Convert string alignment to enum value
                        alignment_str = para_data["alignment"].replace("PP_ALIGN.", "")
                        alignment_value = getattr(PP_ALIGN, alignment_str, None)
                        if alignment_value is not None:
                            para.alignment = alignment_value
                    
                    # Add text runs
                    for run_data in para_data.get("runs", []):
                        run = para.add_run()
                        run.text = run_data["text"]
                        
                        # Apply font formatting
                        font_data = run_data.get("font", {})
                        if font_data.get("name"):
                            run.font.name = font_data["name"]
                        if font_data.get("size"):
                            run.font.size = Pt(font_data["size"])
                        if font_data.get("bold") is not None:
                            run.font.bold = font_data["bold"]
                        if font_data.get("italic") is not None:
                            run.font.italic = font_data["italic"]
                        if font_data.get("underline") is not None:
                            run.font.underline = font_data["underline"]
                        if font_data.get("color"):
                            color_str = font_data["color"].replace("RGB(", "").replace(")", "")
                            rgb_values = [int(x) for x in color_str.split(",")]
                            run.font.color.rgb = RGBColor(*rgb_values)
                
                # Apply text frame properties
                if "has_text_linking" in mod_element and mod_element["has_text_linking"] != orig_element.get("has_text_linking"):
                    shape.text_frame.auto_size = mod_element["has_text_linking"]
                
                if "word_wrap" in mod_element and mod_element["word_wrap"] != orig_element.get("word_wrap"):
                    shape.text_frame.word_wrap = mod_element["word_wrap"]
                
                if "vertical_anchor" in mod_element and mod_element["vertical_anchor"] != orig_element.get("vertical_anchor"):
                    # Convert string anchor to enum value
                    anchor_str = mod_element["vertical_anchor"].replace("MSO_ANCHOR.", "")
                    anchor_value = getattr(MSO_ANCHOR, anchor_str, None)
                    if anchor_value is not None:
                        shape.text_frame.vertical_anchor = anchor_value
        
        # Handle table content changes
        elif mod_element.get("content_type") == "table" and hasattr(shape, "table"):
            mod_table = mod_element.get("table_data", [])
            orig_table = orig_element.get("table_data", [])
            
            # Only update if there are differences
            if mod_table != orig_table:
                for row_idx, row_data in enumerate(mod_table):
                    for col_idx, cell_data in enumerate(row_data):
                        # Check if cell data has changed
                        orig_cell_data = next(
                            (c for c in orig_table[row_idx] if c["row_idx"] == row_idx and c["col_idx"] == col_idx),
                            None
                        )
                        
                        if orig_cell_data and cell_data["text"] != orig_cell_data["text"]:
                            cell = shape.table.cell(row_idx, col_idx)
                            cell.text = cell_data["text"]
        
        # Handle chart content changes
        elif mod_element.get("content_type") == "chart" and hasattr(shape, "chart"):
            # Check if chart title has changed
            if "title" in mod_element and "title" in orig_element:
                if mod_element["title"] != orig_element["title"] and shape.chart.has_title:
                    shape.chart.chart_title.text_frame.text = mod_element["title"]
            
            # Updating chart data is more complex and may require reconstructing the chart
            # This is a simplified version - for production code you'd need more extensive logic
            if "series" in mod_element and "series" in orig_element:
                if mod_element["series"] != orig_element["series"]:
                    print("Warning: Chart data updates require more complex handling")
                    # Would need to update chart data here
        
        # Handle picture content changes (limited capabilities)
        elif mod_element.get("content_type") == "picture" and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Images typically need to be replaced entirely rather than modified
            pass
        
        # Handle shape fill changes
        if "fill_color" in mod_element and "fill_color" in orig_element:
            if mod_element["fill_color"] != orig_element["fill_color"]:
                color_str = mod_element["fill_color"].replace("RGB(", "").replace(")", "")
                rgb_values = [int(x) for x in color_str.split(",")]
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb_values)
    
    return slide