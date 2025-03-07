from datetime import datetime
from http.client import HTTPException
import os
import shutil
from typing import List, Tuple
from pptx import Presentation
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER, MSO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.shapes.placeholder import SlidePlaceholder
from app.schemas import schemas
from app.models.models import SlideMetadata
from app.utils.openai import get_completion, get_formatted_completion, get_embedding
from sqlalchemy.orm import Session
from collections import defaultdict
import json
import numpy as np
import logging
import io
from pathlib import Path

logger = logging.getLogger(__name__)

async def generate_presentation_outline(input_data: schemas.PresentationInput) -> List[schemas.SlideOutline]:
    """Generate a structured outline based on presentation input"""
    system_prompt = """You are an expert presentation designer. Generate a presentation outline in a structured format.
    The response should be a JSON object with a 'slides' array containing the outline sections.
    
    Example output for a presentation with "Number of Slides: 5":
    {
        "slides": [
            {
                "slide_number": 1,
                "section": "introduction",
                "description": "Company introduction with focus on enterprise software expertise",
                "keywords": ["company", "introduction", "expertise", "enterprise"]
            },
            {
                "slide_number": 2,
                "section": "problem",
                "description": "Common challenges faced by enterprises with legacy systems integration",
                "keywords": ["challenges", "pain points", "legacy", "integration"]
            },
            {
                "slide_number": 3,
                "section": "solution",
                "description": "Our approach to modernizing enterprise applications while maintaining business continuity",
                "keywords": ["solution", "modernization", "continuity", "approach"]
            },
            {
                "slide_number": 4,
                "section": "timeline",
                "description": "Proposed implementation timeline for client project",
                "keywords": ["timeline", "implementation", "schedule", "phases"]
            },
            {
                "slide_number": 5,
                "section": "next_steps",
                "description": "Recommended action items to move forward with the engagement",
                "keywords": ["next steps", "action", "engagement", "proposal"]
            }
        ]
    }
    """
    
    user_prompt = f"""Create a presentation outline for:
    Title: {input_data.title}
    Client: {input_data.client_name}
    Industry: {input_data.industry}
    Description: {input_data.description}
    Target Audience: {input_data.target_audience}
    Key Messages: {', '.join(input_data.key_messages)}
    Number of Slides: {input_data.num_slides or '8-12'}
    Tone: {input_data.tone or 'Professional'}
    Additional Context: {input_data.additional_context or 'N/A'}"""

    response = await get_formatted_completion(
        system_prompt=system_prompt,
        user_prompt=user_prompt,
        format_model=schemas.PresentationOutlineResponse
    )
    
    # Convert to SlideOutline objects with validation
    return [schemas.SlideOutline(**slide.model_dump()) for slide in response.slides]

async def find_matching_slides_remix(
    outline: List[schemas.SlideOutline], 
    db: Session,
    similarity_threshold: float = 0.7
) -> List[Tuple[SlideMetadata, int]]:
    """Find the best matching slides for each outline section using semantic search"""
    matched_slides = []
    used_slide_ids = set()  # Track already matched slide IDs
    
    for section in outline:
        search_content = {
            "section": section.section,
            "description": section.description,
            "keywords": section.keywords
        }
        search_embedding = await get_embedding(json.dumps(search_content))
        
        slides = db.query(SlideMetadata).all()
        best_match = None
        best_similarity = -1
        
        for slide in slides:
            if slide.embedding is None or slide.id in used_slide_ids:
                continue  # Skip already selected slides
                
            similarity = np.dot(search_embedding, slide.embedding) / (
                np.linalg.norm(search_embedding) * np.linalg.norm(slide.embedding)
            )

            score_multiplier = 1.0
            if section.section.lower() == slide.category.lower():
                score_multiplier *= 1.2

            final_similarity = similarity * score_multiplier

            if final_similarity > best_similarity:
                best_similarity = final_similarity
                best_match = (slide, slide.id)
        
        if best_match and best_similarity >= similarity_threshold:
            matched_slides.append(best_match)
            used_slide_ids.add(best_match[1])  # Add slide.id to the used set
        else:
            print(f"Warning: No good match found for section {section.section}")
    
    return matched_slides


def modify_ppt_text(file_path, replacements, output_path):
    """
    Modifies text in a PowerPoint file while preserving formatting and handling multi-run text issues.

    :param file_path: Path to the original PowerPoint file.
    :param replacements: Dictionary mapping old text to new text.
    :param output_path: Path to save the modified PowerPoint.
    :return: Success or error message.
    """
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    prs = Presentation(file_path)
    modified = False  # Track if changes were made

    # merger everything we got back into a single dictionary
    merged_replacements = {}

    for item in replacements:
        content = item.get("content", {})
        if "original" in content and "modified" in content:
            merged_replacements[content["original"]] = content["modified"]

    print('--------------------------------------------')
    print(f"replacements: {merged_replacements}")
    print('--------------------------------------------')
    
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
        
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    full_text = paragraph.text.strip()  # Extract full text from paragraph

                    # Check if the entire text matches a replacement
                    if full_text in merged_replacements:
                        new_text = merged_replacements[full_text]  # Get new text
                        runs = paragraph.runs  # Get original text runs                      

                        if len(runs) == 1:
                            print(f"Replacing text: {runs[0].text} with -> {new_text}")
                            # Simple case: Only one run
                            runs[0].text = new_text
                        else:
                            # Complex case: Multiple runs
                            first_run = runs[0]  # Store first run for formatting reference

                            # Clear all runs except the first one
                            print("more complex replacement with multiple runs")
                            print('clear all runs except the first one')
                            for run in runs:
                                print(f"Run Text Before: {run.text}")
                                run.text = ""

                            # Apply new text while keeping first run’s formatting
                            first_run.text = new_text
                            print(f'applied the new text {new_text}')
                            print('while keeping first run’s formatting')
                            
                        modified = True 

    if modified:
        prs.save(output_path)
        return {"message": f"Modified PowerPoint saved as: {output_path}"}
    else:
        prs.save(output_path)
        return {"message": "No matching text found to replace."}
    

def getOriginals(file_path, slide_ids):
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    prs = Presentation(file_path)
    originals = []  # Track if changes were made

    for slide_index, slide in enumerate(prs.slides, start=1):
        if slide_index in slide_ids:
            originals.append(slide)
    return originals


async def generate_slide_content_remix(
    slides: List[SlideMetadata],
    outline: List[schemas.SlideOutline],
    presentation_input: schemas.PresentationInput,
    original_slides_content,
    max_retries: int = 3
) -> List[dict]:
    """Generate customized content for each slide based on the outline"""
    slide_content_list = []
    
    system_prompt = """You are an expert presentation content writer. 
    Your task is to modify only the text content of a PowerPoint slide while preserving its exact structure.
    
    Rules:
    1. Only modify text content where appropriate 
    2. if replacement text is too long, shorten it to about the same number of characters as the original
    3. Keep text concise and impactful
    4. return the modified text as the original in a key value pair like {"original":"what you generated"}
    5. Please only use the same slide once
    """
    grouped_slides = defaultdict(list)

    for slide_id, shape_type, text_content in original_slides_content:
        grouped_slides[slide_id].append({
            "shape_type": shape_type,
            "text_content": text_content
        })
        
    for slide, section in zip(slides, outline):
        for attempt in range(max_retries):
            try:
                prompt = f"""Generate new content for this slide with the following context:
                Slide Type: {slide.slide_type}
                Purpose: {section.description}
                Target Audience: {presentation_input.target_audience}
                Tone: {presentation_input.tone or 'Professional'}
                Slide Number: {slide.slide_number}
                
                Original slide content:
                {grouped_slides}
                
                Return only the original text and the modified text as a key value pair in a JSON object, no other text.
                The Slide Number corresponds to first item in the original slide content tuples
                """
                
                modified_content = await get_completion(
                    prompt=prompt,
                    system_prompt=system_prompt
                )
                
                logger.debug(f"OpenAI response for slide {slide.id} (attempt {attempt + 1}):\n{modified_content}")
                
                # Parse the response as JSON
                content_dict = json.loads(modified_content)
                slide_content_list.append({
                    "slide_id": slide.id,
                    "content": content_dict
                })
                break  # Success, exit retry loop
                
            except json.JSONDecodeError as e:
                logger.warning(f"Attempt {attempt + 1} failed for slide {slide.id}: {e}")
                logger.warning(f"Raw response:\n{modified_content}")
                
                if attempt == max_retries - 1:  # Last attempt failed
                    logger.error(f"All {max_retries} attempts failed for slide {slide.id}")
                    raise ValueError(f"Failed to generate valid JSON for slide {slide.id} after {max_retries} attempts")
    
    return slide_content_list


def construct_presentation_remix(original_slides, output_path, slide_data):
    """Copy selected slides from template and modify content"""

    path_to_copy_project = copyOG(original_slides, slide_data)
    return modify_ppt_text(path_to_copy_project, slide_data, output_path) 

def copyOG(original_slides, replacements):

    merged_replacements = defaultdict(dict)
    slide_ids = set()

    for item in replacements:
        slide_id = item.get("slide_id")
        content = item.get("content", {})

        if isinstance(content, dict):
            merged_replacements.update(content)
            slide_ids.add(slide_id)

    slide_ids = sorted(slide_ids)  # Convert to sorted list if needed
    merged_replacements = dict(merged_replacements)  # Convert back to dictionary

    
    output_dir = Path("presentation_output")
    output_dir.mkdir(exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"duplicated_{timestamp}.pptx"
    output_path = output_dir / output_filename
    
    # First, copy the entire file
    shutil.copy2(original_slides, output_path)
    
    # Then open the copied file and modify it
    prs = Presentation(output_path)
    
    
    # Remove slides that aren't in our keep list
    # We need to remove from end to start to avoid index shifting
    zero_based_slide_ids = [i - 1 for i in slide_ids]  # Convert to 0-based indexing

    for i in range(len(prs.slides) - 1, -1, -1): 
        if i not in zero_based_slide_ids:  
            xml_slides = prs.slides._sldIdLst
            xml_slides.remove(xml_slides[i])
    
    prs.save(output_path)
    return str(output_path)