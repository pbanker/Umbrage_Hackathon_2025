from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER, MSO_SHAPE_TYPE
from datetime import datetime
from typing import Union, Tuple, List, Dict, BinaryIO
from sqlalchemy.orm import Session
from app.models.models import SlideMetadata, SlideShape
import shutil
from pathlib import Path
import json
from app.utils.openai import get_embedding
import uuid
from pdf2image import convert_from_path
import tempfile
import subprocess



async def process_powerpoint_repository(
    pptx_source: Union[str, BinaryIO, bytes],
    db: Session,
    source_type: str = "file_path"
) -> Tuple[str, List[SlideMetadata], List[str]]:
    """
    Process a PowerPoint file and create SlideMetadata objects.
    Stores the PowerPoint file in slides_storage directory.
    
    Args:
        pptx_source: Can be either:
            - file path (str)
            - file-like object (BinaryIO from upload)
            - bytes (from MS Graph API)
        db: Session: SQLAlchemy database session
        source_type: One of "file_path", "upload", or "ms_graph"
    
    Returns:
        Tuple containing:
            - storage_path: Path where the presentation was stored
            - List of SlideMetadata objects
            - List of image paths
    """
    # Create storage directory if it doesn't exist
    storage_dir = Path("slides_repository")
    storage_dir.mkdir(exist_ok=True)
    
    # Generate unique filename using timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    storage_path = storage_dir / f"presentation_{timestamp}.pptx"
    
    # Handle different input types and save to storage
    if source_type == "file_path":
        shutil.copy2(pptx_source, storage_path)
        presentation = Presentation(pptx_source)
    elif source_type == "upload":
        with open(storage_path, "wb") as f:
            if hasattr(pptx_source, "seek"):
                pptx_source.seek(0)
            f.write(pptx_source.read())
        presentation = Presentation(storage_path)
    elif source_type == "ms_graph":
        with open(storage_path, "wb") as f:
            f.write(pptx_source)
        presentation = Presentation(storage_path)
    else:
        raise ValueError("Invalid source_type. Must be 'file_path', 'upload', or 'ms_graph'")

    # After saving the file, generate images
    image_paths = _save_slides_as_images(str(storage_path))
    
    # Create metadata objects for each slide
    slide_metadata_objects = []
    
    for slide_idx, (slide, image_path) in enumerate(zip(presentation.slides, image_paths)):
        semantic_content = {
            "title": _extract_slide_title(slide),
            "purpose": _infer_slide_purpose(slide),
            "category": _infer_slide_category(slide),
            "tags": _generate_slide_tags(slide)
        }
        stringified_metadata = json.dumps(semantic_content)
        embedding = await get_embedding(stringified_metadata)
        
        metadata = SlideMetadata(
            title=semantic_content["title"],
            category=semantic_content["category"],
            slide_type=_infer_slide_type(slide),
            purpose=semantic_content["purpose"],
            tags=semantic_content["tags"],
            audience=None,  # To be filled by user/AI later
            sales_stage=None,  # To be filled by user/AI later
            content_mapping=_create_content_mapping(slide),
            embedding=embedding,
            image_path=image_path,  # Add the image path
            slide_number= slide_idx + 1
        )
        
        slide_metadata_objects.append(metadata)
    
    return str(storage_path), slide_metadata_objects, image_paths

def _extract_slide_title(slide) -> str:
    """Extract the title from a slide"""
    for shape in slide.shapes:
        try:
            # First check if it's a placeholder and if it's a title
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                if shape.has_text_frame:
                    return shape.text_frame.text
        except (AttributeError, ValueError):
            # If it's not a placeholder, check if it might be a title shape by name
            if hasattr(shape, 'name') and 'Title' in shape.name:
                if shape.has_text_frame:
                    return shape.text_frame.text
    
    return "Untitled Slide"

def _infer_slide_category(slide) -> str:
    """Infer the category of the slide based on its content"""
    # Check for charts
    if any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in slide.shapes):
        return "data_visualization"
    
    # Check for tables
    if any(hasattr(shape, "table") for shape in slide.shapes):
        return "data_presentation"
    
    # Check for multiple images
    if sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE) > 1:
        return "visual"
    
    # Default to content slide
    return "content"

def _infer_slide_type(slide) -> str:
    """Infer the specific type of slide based on its content and layout"""
    layout_name = slide.slide_layout.name.lower()
    
    # Check layout name first
    if "title" in layout_name:
        if "content" in layout_name:
            return "title_and_content"
        return "title_slide"
    
    # Check content types
    if any(shape.shape_type == MSO_SHAPE_TYPE.CHART for shape in slide.shapes):
        return "chart_slide"
    
    if any(hasattr(shape, "table") for shape in slide.shapes):
        return "table_slide"
    
    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        return "image_slide"
    
    return "text_slide"

def _infer_slide_purpose(slide) -> str:
    """Infer the purpose of the slide based on its type and content"""
    slide_type = _infer_slide_type(slide)
    
    # Basic mapping of slide types to purposes
    purpose_map = {
        "title_slide": "To introduce the presentation",
        "chart_slide": "To visualize data or trends",
        "table_slide": "To present structured data",
        "image_slide": "To provide visual information",
        "text_slide": "To convey textual information"
    }
    
    return purpose_map.get(slide_type, "To present information")

def _generate_slide_tags(slide) -> List[str]:
    """Generate relevant tags for the slide"""
    tags = []
    
    # Add layout-based tag
    tags.append(slide.slide_layout.name.lower())
    
    # Add content-type based tags
    content_types = set()
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            content_types.add("chart")
        elif hasattr(shape, "table"):
            content_types.add("table")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            content_types.add("picture")
        elif shape.has_text_frame:
            content_types.add("text")
    
    tags.extend(list(content_types))
    
    return tags

#TODO: This is a placeholder function to create a content object/schema for the slides.
def _create_content_mapping(slide) -> Dict:
    """
    Create comprehensive content object for a given slide.
    
    Args:
        slide (Slide): The slide to create a content object for
    
    Returns:
        dict: Dictionary containing the content object for the slide
    """
    # Get placeholder type names for better readability
    placeholder_types = {
        getattr(PP_PLACEHOLDER, attr): attr 
        for attr in dir(PP_PLACEHOLDER) 
        if not attr.startswith('__')
    }
    
    # Get shape type names
    shape_types = {
        getattr(MSO_SHAPE_TYPE, attr): attr 
        for attr in dir(MSO_SHAPE_TYPE) 
        if not attr.startswith('__')
    }
    
    # Basic slide information
    schema = {
        "slide_id": slide.slide_id,
        "layout_name": slide.slide_layout.name,
        "elements": []
    }
    
    # Process all shapes on the slide
    for shape in slide.shapes:
        # Base shape information
        element = {
            "id": shape.shape_id,
            "name": shape.name,
            "shape_type": shape_types.get(shape.shape_type, str(shape.shape_type)),
            "position": {
                "x": shape.left,
                "y": shape.top,
                "width": shape.width,
                "height": shape.height,
                "rotation": shape.rotation
            }
        }
        
        # Add placeholder information if applicable
        try:
            if shape.is_placeholder:
                element["is_placeholder"] = True
                placeholder_type = shape.placeholder_format.type
                element["placeholder_type"] = placeholder_types.get(placeholder_type, str(placeholder_type))
                element["placeholder_idx"] = shape.placeholder_format.idx
            else:
                element["is_placeholder"] = False
        except (AttributeError, ValueError):
            element["is_placeholder"] = False
        
        # Text content
        if shape.has_text_frame:
            element["content_type"] = "text"
            
            # Get text and formatting at paragraph level
            paragraphs_data = []
            for p in shape.text_frame.paragraphs:
                paragraph_data = {
                    "text": p.text,
                    "level": p.level,
                    "alignment": str(p.alignment) if hasattr(p, 'alignment') else None,
                    "runs": []
                }
                
                # Get formatting for each text run
                for run in p.runs:
                    run_data = {
                        "text": run.text,
                        "font": {
                            "name": run.font.name,
                            "size": run.font.size.pt if hasattr(run.font.size, 'pt') else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                            "color": str(run.font.color.rgb) if hasattr(run.font.color, 'rgb') and run.font.color.rgb else None
                        }
                    }
                    paragraph_data["runs"].append(run_data)
                
                paragraphs_data.append(paragraph_data)
            
            element["paragraphs"] = paragraphs_data
            element["has_text_linking"] = shape.text_frame.auto_size
            element["word_wrap"] = shape.text_frame.word_wrap
            element["vertical_anchor"] = str(shape.text_frame.vertical_anchor) if hasattr(shape.text_frame, 'vertical_anchor') else None
        
        # Table content
        elif hasattr(shape, 'table'):
            element["content_type"] = "table"
            table_data = []
            
            for r_idx, row in enumerate(shape.table.rows):
                row_data = []
                for c_idx, cell in enumerate(row.cells):
                    cell_data = {
                        "text": cell.text,
                        "row_idx": r_idx,
                        "col_idx": c_idx,
                        "width": shape.table.columns[c_idx].width,
                        "height": shape.table.rows[r_idx].height
                    }
                    row_data.append(cell_data)
                table_data.append(row_data)
            
            element["table_data"] = table_data
            element["row_count"] = len(shape.table.rows)
            element["column_count"] = len(shape.table.columns)
        
        # Chart content
        elif hasattr(shape, 'chart'):
            element["content_type"] = "chart"
            element["chart_type"] = str(shape.chart.chart_type)
            
            # Extract categories and series
            if hasattr(shape.chart, 'plots') and shape.chart.plots:
                plot = shape.chart.plots[0]
                
                # Try to get categories
                categories = []
                if hasattr(plot, 'categories') and plot.categories:
                    for category in plot.categories:
                        categories.append(str(category.label))
                element["categories"] = categories
                
                # Try to get series
                series_data = []
                if hasattr(plot, 'series'):
                    for series in plot.series:
                        series_info = {
                            "name": series.name if hasattr(series, 'name') else "Unknown",
                            "values": list(series.values) if hasattr(series, 'values') else []
                        }
                        series_data.append(series_info)
                element["series"] = series_data
            
            element["has_title"] = shape.chart.has_title
            if shape.chart.has_title:
                element["title"] = shape.chart.chart_title.text_frame.text
        
        # Picture content
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["content_type"] = "picture"
            if hasattr(shape, 'image'):
                element["image_filename"] = shape.image.filename if hasattr(shape.image, 'filename') else None
                element["image_format"] = shape.image.ext if hasattr(shape.image, 'ext') else None
                element["content_type"] = f"image/{shape.image.ext}" if hasattr(shape.image, 'ext') else "image"
        
        # Shape content
        else:
            element["content_type"] = "shape"
            if hasattr(shape, 'fill'):
                element["fill_type"] = str(shape.fill.type) if hasattr(shape.fill, 'type') else None
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                    element["fill_color"] = str(shape.fill.fore_color.rgb) if hasattr(shape.fill.fore_color, 'rgb') else None
        
        # Add this element to the slide schema
        schema["elements"].append(element)
    
    # Add background information
    try:
        if hasattr(slide, 'background') and slide.background:
            if hasattr(slide.background, 'fill'):
                fill_type = type(slide.background.fill._fill).__name__
                schema["background_fill_type"] = fill_type
                
                # Only try to get fore_color if it's a solid fill
                if fill_type == "SolidFill" and hasattr(slide.background.fill, 'fore_color'):
                    if slide.background.fill.fore_color and hasattr(slide.background.fill.fore_color, 'rgb'):
                        schema["background_color"] = str(slide.background.fill.fore_color.rgb)
    except (AttributeError, TypeError):
        # If anything goes wrong with background processing, just skip it
        pass
    
    return schema

def _save_slides_as_images(pptx_path: str) -> list[str]:
    """
    Convert each slide in the PowerPoint to an image and save it.
    Uses LibreOffice to convert to PDF first, then pdf2image to convert to images.
    """
    # Create images directory if it doesn't exist
    images_dir = Path("images")
    images_dir.mkdir(exist_ok=True)
    
    # Create a temporary directory for the PDF
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        pdf_path = temp_dir_path / "slides.pdf"
        
        # Convert PPTX to PDF using LibreOffice
        try:
            # Get absolute paths
            pptx_abs_path = str(Path(pptx_path).absolute())
            temp_dir_abs_path = str(temp_dir_path.absolute())
            
            # Try using libreoffice with absolute paths
            result = subprocess.run([
                'libreoffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', temp_dir_abs_path,
                pptx_abs_path
            ], capture_output=True, text=True, check=False)
            
            if result.returncode != 0:
                print(f"LibreOffice Error: {result.stderr}")
                # Try soffice as fallback
                result = subprocess.run([
                    'soffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir_abs_path,
                    pptx_abs_path
                ], capture_output=True, text=True, check=False)
                
                if result.returncode != 0:
                    print(f"Soffice Error: {result.stderr}")
                    raise RuntimeError(
                        f"Failed to convert PPTX to PDF. LibreOffice/Soffice error: {result.stderr}"
                    )
            
            # Verify PDF was created
            expected_pdf = temp_dir_path / Path(pptx_path).with_suffix('.pdf').name
            if not expected_pdf.exists():
                raise RuntimeError(f"PDF file not created at expected location: {expected_pdf}")
            
            # Move the PDF to the expected location
            shutil.move(str(expected_pdf), str(pdf_path))
            
            # Convert PDF to images
            images = convert_from_path(str(pdf_path))
            image_paths = []
            
            # Save each image
            for i, image in enumerate(images):
                image_filename = f"slide_{uuid.uuid4()}.jpg"
                image_path = images_dir / image_filename
                image.save(str(image_path), "JPEG")
                image_paths.append(str(image_path))
            
            return image_paths
            
        except Exception as e:
            print(f"Error during conversion: {str(e)}")
            raise RuntimeError(f"Failed to process slides: {str(e)}")
        

def retrieve_shape_and_content(pptx_storage_path: str):
  
    prs = Presentation(pptx_storage_path)

    shapes_to_insert = []  # all my shapes
    for slide_index, slide in enumerate(prs.slides, start=1):
      for shape_index, shape in enumerate(slide.shapes, start=1):
      
          shape_type_id = shape.shape_type  # Gets shape type ID
          shape_type_name = MSO_SHAPE_TYPE(shape_type_id).name if shape_type_id in MSO_SHAPE_TYPE.__members__.values() else f"Unknown ({shape_type_id})"

          if hasattr(shape, "text_frame") and shape.text_frame is not None:
              for paragraph in shape.text_frame.paragraphs:
                  full_text = paragraph.text.strip()  # Extract full text from paragraph
                  print(f"Slide {slide_index}, Shape {shape_index} [{shape_type_name}]: {full_text}")
                  new_shape = SlideShape(
                      slide_metadata_id=slide_index,
                      shape_index=shape_index,
                      shape_type=shape_type_name,
                      text_content=full_text
                     )
                  shapes_to_insert.append(new_shape)
    if shapes_to_insert:
      return shapes_to_insert
    else:
        return None